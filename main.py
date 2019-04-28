from pptx import Presentation
from pptx.util import Pt
import re


with open('input.txt') as f:
    text = f.read()
    text_splitted = text.strip().split('\n')
i = 8 + 7 
blocks = []
main_file_name = text_splitted[0]
main_title = text_splitted[1]
for result in re.findall('---(.*?)---', text, re.S):
    title = re.findall('\*\*(.+)', result)[0]
    bullets = re.findall('-(.+)', result)
    blocks.append({'title': title, 'bullets': bullets})


def main(template_name, main_file_name, main_title, blocks):
    prs = Presentation('templates/' + template_name)

    # Editing the first slide
    slide = prs.slides[0]
    slide.placeholders[0].text = main_title

    bullet_slide_layout = prs.slide_layouts[2]

    for block in blocks:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = block['title'].strip()

        tf = body_shape.text_frame
        k = 0
        points = len(block['bullets'])
        total = points * 2
        for j in range(total):
            p = tf.paragraphs[j]
            run = p.add_run()
            if j%2 != 0:
                run.text = block['bullets'][k].strip()
                k+=1
            else:
                run.text = ''
            font = run.font
            font.size = Pt(25 - (points * 1.7))
            tf.add_paragraph()
            tf.add_paragraph()


    prs.save('outputs/' + main_file_name + '.pptx')

main('template-normal.pptx', main_file_name, main_title, blocks)
main('template-online.pptx', main_file_name + ' - online', main_title, blocks)